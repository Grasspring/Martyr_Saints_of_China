import csv
from pptx import Presentation

def pptx_to_csv(pptx_file, csv_file):
    # Load the presentation
    prs = Presentation(pptx_file)
    
    # Open the CSV file for writing
    with open(csv_file, mode='w', newline='', encoding='utf-8') as file:
        writer = csv.writer(file)
        
        # Write the header
        header = ['Slide number'] + [f'text {i+1}' for i in range(4)]
        writer.writerow(header)
        
        # Iterate through each slide
        for slide_number, slide in enumerate(prs.slides, start=1):
            texts = []
            
            # Collect text from each shape on the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text.strip())
            
            # Limit to the first 4 texts or fill with empty strings
            texts = texts[:4] + [''] * (4 - len(texts))
            
            # Write the slide number and texts to the CSV
            writer.writerow([slide_number] + texts)

# Example usage
pptx_file = 'final_result.pptx'  # Replace with your pptx file path
csv_file = 'result_output.csv'  # Output CSV file path
pptx_to_csv(pptx_file, csv_file)