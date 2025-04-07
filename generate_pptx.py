import pandas as pd
from pptx import Presentation

# Load the CSV file
df = pd.read_csv('data.csv')

# Load the template presentation
presentation = Presentation('template.pptx')

# Loop through each row in the DataFrame
for index, row in df.iterrows():
    # Copy the first slide
    print(f"Processing character: {row['Name']}")
    slide = presentation.slides[index]
    run = slide.shapes.title.text_frame.paragraphs[0].add_run()
    run.text = row['Name']  # Set the title to the character's name
    slide.placeholders[1].text = row['Story']  # Set the story text
    image = row['Image']  # Get the image path
    image_path = f"images/{image}"  # Construct the image path
   
# Save the new presentation
presentation.save('output_presentation.pptx')